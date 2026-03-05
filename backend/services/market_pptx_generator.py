"""
Market Study PPTX Generator: creates a 10-slide market study presentation.

Programmatically builds slides using python-pptx with Constellation Capital AG
brand colors. Each slide maps to one section of MarketStudyData.
"""

import io
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from models.market_study import MarketStudyData

# Brand colors
DARK_NAVY = RGBColor(0x22, 0x3F, 0x6A)
DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
MID_BLUE = RGBColor(0x2E, 0x75, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
GREEN = RGBColor(0x00, 0xB0, 0x50)
RED = RGBColor(0xC0, 0x00, 0x00)
GOLD = RGBColor(0xF4, 0xC5, 0x00)
MEDIUM_GRAY = RGBColor(0x80, 0x80, 0x80)

# Slide dimensions (widescreen 10" x 5.625")
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)


def _add_run(paragraph, text, font_name="Calibri", font_size=Pt(9),
             bold=False, color=DARK_GRAY, italic=False):
    """Add a formatted run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    if italic:
        run.font.italic = True
    return run


def _add_title_bar(slide, title_text: str) -> None:
    """Add a dark navy title bar at the top of the slide."""
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.7),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_NAVY
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, title_text, font_size=Pt(18), bold=True, color=WHITE)


def _add_subtitle(slide, text: str, top: float = 0.8) -> None:
    """Add a subtitle below the title bar."""
    txbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(top), Inches(9), Inches(0.4),
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_run(p, text, font_size=Pt(11), color=MID_BLUE, italic=True)


def _add_bullet_list(slide, items: list[str], left: float, top: float,
                     width: float, height: float, font_size=Pt(9)) -> None:
    """Add a bullet-point list to the slide."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        _add_run(p, f"• {item}", font_size=font_size, color=DARK_GRAY)


def _add_table(slide, headers: list[str], rows: list[list[str]],
               left: float, top: float, width: float, row_height: float = 0.3) -> None:
    """Add a styled table to the slide."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(row_height * n_rows),
    )
    table = table_shape.table

    # Style header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(8)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY

    # Style data rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(0xF2, 0xF2, 0xF2)


def _rating_color(rating: str) -> RGBColor:
    """Map a rating string to a color."""
    r = rating.lower()
    if r in ("positive", "low", "high"):
        return GREEN
    if r in ("negative",):
        return RED
    return GOLD


# ─── Slide builders ────────────────────────────────────────────────────────

def _slide_1_executive_summary(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 1: Executive Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_title_bar(slide, data.executive_summary.title or f"Market Study: {data.meta.market_name}")

    if data.executive_summary.market_verdict:
        _add_subtitle(slide, data.executive_summary.market_verdict)

    if data.executive_summary.key_findings:
        _add_bullet_list(
            slide, data.executive_summary.key_findings,
            left=0.5, top=1.4, width=9.0, height=3.5, font_size=Pt(11),
        )

    # Footer with meta
    txbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.4),
    )
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    meta_text = f"{data.meta.market_name} | {data.meta.region} | {data.meta.research_date}"
    _add_run(p, meta_text, font_size=Pt(7), color=MEDIUM_GRAY, italic=True)


def _slide_2_market_sizing(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 2: Market Sizing (TAM/SAM/SOM)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Market Sizing")

    ms = data.market_sizing

    # TAM / SAM / SOM boxes
    metrics = [
        ("TAM", ms.tam, ms.tam_year),
        ("SAM", ms.sam, ms.sam_year),
        ("SOM", ms.som, ""),
    ]
    for i, (label, value, year) in enumerate(metrics):
        left = 0.5 + i * 3.1
        box = slide.shapes.add_shape(
            1, Inches(left), Inches(0.9), Inches(2.8), Inches(1.0),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_run(p, label, font_size=Pt(10), bold=True, color=MID_BLUE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _add_run(p2, value or "n/a", font_size=Pt(16), bold=True, color=DARK_NAVY)
        if year:
            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.CENTER
            _add_run(p3, year, font_size=Pt(8), color=MEDIUM_GRAY)

    # CAGR
    if ms.cagr is not None:
        txbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.1), Inches(4.0), Inches(0.4),
        )
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        cagr_pct = f"{ms.cagr * 100:.1f}%"
        _add_run(p, f"CAGR: {cagr_pct} ({ms.cagr_period})", font_size=Pt(11), bold=True, color=GREEN)

    # Data points table
    if ms.data_points:
        headers = ["Year", "Value", "Label"]
        rows = [
            [dp.year, f"EUR {dp.value:.1f}bn" if dp.value else "n/a", dp.label]
            for dp in ms.data_points
        ]
        _add_table(slide, headers, rows, left=0.5, top=2.7, width=5.0)

    # Assumptions
    if ms.assumptions:
        _add_bullet_list(
            slide, ms.assumptions,
            left=5.8, top=2.7, width=3.7, height=2.5, font_size=Pt(8),
        )


def _slide_3_segmentation(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 3: Market Segmentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Market Segmentation")

    if data.market_segments:
        headers = ["Segment", "Size", "Share", "Growth", "Description"]
        rows = [
            [
                seg.name,
                seg.size or "n/a",
                f"{seg.share_pct:.0f}%" if seg.share_pct is not None else "n/a",
                seg.growth_rate or "n/a",
                seg.description[:60] + "..." if len(seg.description) > 60 else seg.description,
            ]
            for seg in data.market_segments
        ]
        _add_table(slide, headers, rows, left=0.3, top=0.9, width=9.4)
    else:
        _add_subtitle(slide, "No segment data available", top=2.0)


def _slide_4_competitive_landscape(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 4: Competitive Landscape."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Competitive Landscape")

    cl = data.competitive_landscape

    # Fragmentation indicator
    txbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.9), Inches(4.0), Inches(0.4),
    )
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    _add_run(p, f"Market Fragmentation: ", font_size=Pt(10), color=DARK_GRAY)
    frag_color = GREEN if cl.fragmentation == "high" else (GOLD if cl.fragmentation == "medium" else RED)
    _add_run(p, cl.fragmentation.upper(), font_size=Pt(10), bold=True, color=frag_color)

    if cl.avg_company_revenue:
        p2 = tf.add_paragraph()
        _add_run(p2, f"Avg. Company Revenue: {cl.avg_company_revenue}", font_size=Pt(9), color=MEDIUM_GRAY)

    # Competitor table
    if cl.top_players:
        headers = ["Company", "Market Share", "Revenue", "HQ", "Key Strengths"]
        rows = [
            [
                cp.name,
                cp.market_share or "n/a",
                cp.revenue or "n/a",
                cp.hq or "n/a",
                "; ".join(cp.strengths[:2]) if cp.strengths else "n/a",
            ]
            for cp in cl.top_players[:7]
        ]
        _add_table(slide, headers, rows, left=0.3, top=1.7, width=9.4)

    # Consolidation trend
    if cl.consolidation_trend:
        txbox2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.8), Inches(9.0), Inches(0.4),
        )
        tf2 = txbox2.text_frame
        p = tf2.paragraphs[0]
        _add_run(p, f"Consolidation: {cl.consolidation_trend}", font_size=Pt(8), color=MEDIUM_GRAY, italic=True)


def _slide_5_trends(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 5: Trends & Drivers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Market Trends & Drivers")

    td = data.trends_drivers

    # Left column: Growth Drivers
    if td.growth_drivers:
        txbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.9), Inches(4.5), Inches(0.3))
        tf = txbox.text_frame
        _add_run(tf.paragraphs[0], "Growth Drivers", font_size=Pt(11), bold=True, color=GREEN)
        _add_bullet_list(slide, td.growth_drivers, left=0.3, top=1.3, width=4.5, height=1.8)

    # Right column: Headwinds
    if td.headwinds:
        txbox = slide.shapes.add_textbox(Inches(5.2), Inches(0.9), Inches(4.5), Inches(0.3))
        tf = txbox.text_frame
        _add_run(tf.paragraphs[0], "Headwinds", font_size=Pt(11), bold=True, color=RED)
        _add_bullet_list(slide, td.headwinds, left=5.2, top=1.3, width=4.5, height=1.8)

    # Bottom left: Tech shifts
    if td.technological_shifts:
        txbox = slide.shapes.add_textbox(Inches(0.3), Inches(3.2), Inches(4.5), Inches(0.3))
        tf = txbox.text_frame
        _add_run(tf.paragraphs[0], "Technological Shifts", font_size=Pt(11), bold=True, color=MID_BLUE)
        _add_bullet_list(slide, td.technological_shifts, left=0.3, top=3.6, width=4.5, height=1.5)

    # Bottom right: Regulatory
    if td.regulatory_changes:
        txbox = slide.shapes.add_textbox(Inches(5.2), Inches(3.2), Inches(4.5), Inches(0.3))
        tf = txbox.text_frame
        _add_run(tf.paragraphs[0], "Regulatory Changes", font_size=Pt(11), bold=True, color=GOLD)
        _add_bullet_list(slide, td.regulatory_changes, left=5.2, top=3.6, width=4.5, height=1.5)


def _slide_6_pestel(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 6: PESTEL Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "PESTEL Analysis")

    pestel = data.pestel
    dimensions = [
        ("Political", pestel.political),
        ("Economic", pestel.economic),
        ("Social", pestel.social),
        ("Technological", pestel.technological),
        ("Environmental", pestel.environmental),
        ("Legal", pestel.legal),
    ]

    # 3x2 grid
    for idx, (label, factor) in enumerate(dimensions):
        col = idx % 3
        row = idx // 3
        left = 0.3 + col * 3.2
        top = 0.9 + row * 2.2

        # Box background
        box = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(3.0), Inches(2.0),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        box.line.color.rgb = LIGHT_GRAY

        # Title + rating
        txbox = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.1), Inches(2.8), Inches(0.3),
        )
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        _add_run(p, f"{label} ", font_size=Pt(9), bold=True, color=DARK_NAVY)
        _add_run(p, f"[{factor.rating}]", font_size=Pt(8), bold=True, color=_rating_color(factor.rating))

        # Points
        if factor.points:
            for i, point in enumerate(factor.points[:3]):
                txbox2 = slide.shapes.add_textbox(
                    Inches(left + 0.1), Inches(top + 0.45 + i * 0.45),
                    Inches(2.8), Inches(0.4),
                )
                tf2 = txbox2.text_frame
                tf2.word_wrap = True
                _add_run(tf2.paragraphs[0], f"• {point}", font_size=Pt(7), color=DARK_GRAY)


def _slide_7_porters(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 7: Porter's Five Forces."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Porter's Five Forces")

    pf = data.porters_five_forces
    forces = [
        ("Competitive Rivalry", pf.rivalry),
        ("Buyer Power", pf.buyer_power),
        ("Supplier Power", pf.supplier_power),
        ("Threat of New Entrants", pf.threat_new_entrants),
        ("Threat of Substitutes", pf.threat_substitutes),
    ]

    for idx, (label, force) in enumerate(forces):
        top = 0.9 + idx * 0.9

        # Rating badge
        badge = slide.shapes.add_shape(
            1, Inches(0.5), Inches(top), Inches(1.2), Inches(0.5),
        )
        badge.fill.solid()
        rating_str = force.rating.upper()
        if rating_str == "HIGH":
            badge.fill.fore_color.rgb = RGBColor(0xFF, 0xE0, 0xE0)
        elif rating_str == "MEDIUM":
            badge.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xD0)
        else:
            badge.fill.fore_color.rgb = RGBColor(0xE0, 0xFF, 0xE0)
        badge.line.fill.background()

        tf = badge.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_run(p, rating_str, font_size=Pt(9), bold=True, color=DARK_GRAY)

        # Label + explanation
        txbox = slide.shapes.add_textbox(
            Inches(1.9), Inches(top), Inches(7.5), Inches(0.7),
        )
        tf2 = txbox.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        _add_run(p2, f"{label}: ", font_size=Pt(10), bold=True, color=DARK_NAVY)
        _add_run(p2, force.explanation, font_size=Pt(9), color=DARK_GRAY)


def _slide_8_value_chain(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 8: Value Chain & Business Models."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Value Chain & Business Models")

    vc = data.value_chain

    # Value chain stages as horizontal flow
    if vc.stages:
        n = len(vc.stages)
        stage_width = min(1.8, 9.0 / max(n, 1))
        for idx, stage in enumerate(vc.stages[:5]):  # Max 5 stages
            left = 0.3 + idx * (stage_width + 0.1)
            box = slide.shapes.add_shape(
                1, Inches(left), Inches(0.9), Inches(stage_width), Inches(1.4),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
            box.line.color.rgb = MID_BLUE

            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _add_run(p, stage.name, font_size=Pt(8), bold=True, color=DARK_NAVY)
            if stage.typical_margin:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                _add_run(p2, stage.typical_margin, font_size=Pt(7), color=GREEN)
            if stage.description:
                p3 = tf.add_paragraph()
                p3.alignment = PP_ALIGN.CENTER
                _add_run(p3, stage.description[:50], font_size=Pt(7), color=MEDIUM_GRAY)

    # Dominant business models
    if vc.dominant_business_models:
        txbox = slide.shapes.add_textbox(Inches(0.3), Inches(2.6), Inches(9.0), Inches(0.3))
        tf = txbox.text_frame
        _add_run(tf.paragraphs[0], "Dominant Business Models", font_size=Pt(11), bold=True, color=DARK_NAVY)
        _add_bullet_list(slide, vc.dominant_business_models, left=0.3, top=3.0, width=9.0, height=1.5)

    # Margin distribution
    if vc.margin_distribution:
        txbox2 = slide.shapes.add_textbox(Inches(0.3), Inches(4.7), Inches(9.0), Inches(0.5))
        tf2 = txbox2.text_frame
        tf2.word_wrap = True
        _add_run(tf2.paragraphs[0], f"Margin Distribution: {vc.margin_distribution}",
                 font_size=Pt(9), color=MEDIUM_GRAY, italic=True)


def _slide_9_buy_and_build(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 9: Buy & Build Potential."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Buy & Build Potential")

    bb = data.buy_and_build

    # Fragmentation score gauge
    if bb.fragmentation_score is not None:
        score_pct = f"{bb.fragmentation_score * 100:.0f}%"
        box = slide.shapes.add_shape(
            1, Inches(0.5), Inches(0.9), Inches(2.5), Inches(0.8),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
        box.line.fill.background()
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_run(p, "Fragmentation Score", font_size=Pt(9), color=MID_BLUE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _add_run(p2, score_pct, font_size=Pt(20), bold=True, color=DARK_NAVY)

    # Estimated targets
    if bb.estimated_targets_dach:
        txbox = slide.shapes.add_textbox(Inches(3.5), Inches(0.9), Inches(6.0), Inches(0.4))
        tf = txbox.text_frame
        _add_run(tf.paragraphs[0], f"Estimated Targets DACH: {bb.estimated_targets_dach}",
                 font_size=Pt(10), bold=True, color=DARK_NAVY)

    # Platform candidates
    if bb.platform_candidates:
        txbox = slide.shapes.add_textbox(Inches(0.3), Inches(2.0), Inches(4.5), Inches(0.3))
        tf = txbox.text_frame
        _add_run(tf.paragraphs[0], "Platform Candidate Profiles", font_size=Pt(10), bold=True, color=DARK_NAVY)
        _add_bullet_list(slide, bb.platform_candidates, left=0.3, top=2.4, width=4.5, height=2.5)

    # Add-on profile + rationale
    right_items = []
    if bb.add_on_profile:
        right_items.append(f"Add-on Profile: {bb.add_on_profile}")
    if bb.consolidation_rationale:
        right_items.append(f"Rationale: {bb.consolidation_rationale}")

    if right_items:
        txbox2 = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.5), Inches(0.3))
        tf2 = txbox2.text_frame
        _add_run(tf2.paragraphs[0], "Consolidation Logic", font_size=Pt(10), bold=True, color=DARK_NAVY)
        _add_bullet_list(slide, right_items, left=5.2, top=2.4, width=4.5, height=2.5, font_size=Pt(8))


def _slide_10_strategic_implications(prs: Presentation, data: MarketStudyData) -> None:
    """Slide 10: Strategic Implications & Recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Strategic Implications & Recommendations")

    si = data.strategic_implications

    # Investment attractiveness
    if si.investment_attractiveness:
        txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9.0), Inches(0.4))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        _add_run(p, "Investment Attractiveness: ", font_size=Pt(11), color=DARK_GRAY)
        attr_color = GREEN if si.investment_attractiveness == "high" else (
            GOLD if si.investment_attractiveness == "medium" else RED
        )
        _add_run(p, si.investment_attractiveness.upper(), font_size=Pt(11), bold=True, color=attr_color)

    # Recommendations (3 boxes)
    for idx, rec in enumerate(si.recommendations[:3]):
        top = 1.5 + idx * 1.2
        box = slide.shapes.add_shape(
            1, Inches(0.5), Inches(top), Inches(9.0), Inches(1.0),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        box.line.color.rgb = LIGHT_GRAY

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_run(p, f"{idx + 1}. {rec.title}", font_size=Pt(10), bold=True, color=DARK_NAVY)
        if rec.description:
            p2 = tf.add_paragraph()
            _add_run(p2, rec.description, font_size=Pt(8), color=DARK_GRAY)
        if rec.risk_benefit:
            p3 = tf.add_paragraph()
            _add_run(p3, f"Risk/Benefit: {rec.risk_benefit}", font_size=Pt(8), color=MID_BLUE, italic=True)

    # Key risks
    if si.key_risks:
        txbox = slide.shapes.add_textbox(Inches(0.3), Inches(4.8), Inches(9.4), Inches(0.3))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        _add_run(p, "Key Risks: ", font_size=Pt(8), bold=True, color=RED)
        _add_run(p, " | ".join(si.key_risks[:4]), font_size=Pt(8), color=DARK_GRAY)


# ─── Main entry point ──────────────────────────────────────────────────────

def generate_market_study(data: MarketStudyData) -> bytes:
    """Generate a 10-slide market study PPTX and return the bytes."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _slide_1_executive_summary(prs, data)
    _slide_2_market_sizing(prs, data)
    _slide_3_segmentation(prs, data)
    _slide_4_competitive_landscape(prs, data)
    _slide_5_trends(prs, data)
    _slide_6_pestel(prs, data)
    _slide_7_porters(prs, data)
    _slide_8_value_chain(prs, data)
    _slide_9_buy_and_build(prs, data)
    _slide_10_strategic_implications(prs, data)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
