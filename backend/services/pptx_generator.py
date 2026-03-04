"""
PPTX Generator: fills the One-Pager template with data and chart images.

Uses the Constellation Capital AG original template (10" x 5.625", 4-column
layout) as the base. Named shapes in the template are filled with structured
data from the OnePagerData model.
"""

import io
import os
import shutil
import tempfile
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from models.one_pager import OnePagerData, CriterionStatus
from services.chart_generator import generate_financials_chart, generate_revenue_donut

# Constellation Capital AG brand colors (from original templates)
DARK_NAVY = RGBColor(0x22, 0x3F, 0x6A)
DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
MID_BLUE = RGBColor(0x2E, 0x75, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
GREEN = RGBColor(0x00, 0xB0, 0x50)
RED = RGBColor(0xC0, 0x00, 0x00)
GOLD = RGBColor(0xF4, 0xC5, 0x00)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)

TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "..", "template")
TEMPLATE_PATH = os.path.join(TEMPLATE_DIR, "one_pager_template.pptx")


def _build_shape_map(slide) -> dict:
    """Build name->shape mapping, including shapes inside groups."""
    shape_map = {}
    for shape in slide.shapes:
        shape_map[shape.name] = shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for child in shape.shapes:
                shape_map[child.name] = child
    return shape_map


def _clear_text_frame(tf):
    """Clear all paragraphs from a text frame."""
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    if tf.paragraphs:
        for run in tf.paragraphs[0].runs:
            run.text = ""


def _add_run(paragraph, text, font_name="Calibri", font_size=Pt(8),
             bold=False, color=None):
    """Add a formatted run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return run


def _set_shape_text(shape, text, font_size=Pt(8), bold=False, color=None,
                    alignment=PP_ALIGN.LEFT):
    """Replace all text in a shape with a single styled run."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    _clear_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = alignment
    _add_run(p, text, font_size=font_size, bold=bold, color=color)


def _fill_header(shape_map: dict, data: OnePagerData):
    """Fill the header section (title placeholder + tagline + investment thesis)."""
    # "One Pager | COMPANY NAME" in the title placeholder
    if "header_title" in shape_map:
        shape = shape_map["header_title"]
        if shape.has_text_frame:
            tf = shape.text_frame
            _clear_text_frame(tf)
            p = tf.paragraphs[0]
            _add_run(p, f"One Pager | ", font_size=Pt(12), bold=True)
            _add_run(p, data.header.company_name, font_size=Pt(12), bold=True)

    # Tagline below the title
    if "header_tagline" in shape_map:
        shape = shape_map["header_tagline"]
        if shape.has_text_frame:
            tf = shape.text_frame
            _clear_text_frame(tf)
            p = tf.paragraphs[0]
            _add_run(p, data.header.tagline or "", font_size=Pt(11))

    # Investment thesis bar (full width)
    if "investment_thesis" in shape_map:
        shape = shape_map["investment_thesis"]
        if shape.has_text_frame:
            tf = shape.text_frame
            _clear_text_frame(tf)
            p = tf.paragraphs[0]
            _add_run(p, "Investment-Thesis:  ", font_size=Pt(8), bold=True)
            _add_run(p, data.investment_thesis, font_size=Pt(8))


def _fill_key_facts(shape_map: dict, data: OnePagerData):
    """Fill the Key Facts section (column 1)."""
    shape = shape_map.get("keyfacts_content")
    if not shape or not shape.has_text_frame:
        return

    tf = shape.text_frame
    _clear_text_frame(tf)
    kf = data.key_facts

    fields = [
        ("Founded:", kf.founded),
        ("HQ:", kf.hq),
        ("Website:", kf.website),
        ("Industry:", f"{kf.industry} / {kf.niche}" if kf.niche else kf.industry),
        ("Revenue:", f"{kf.revenue} ({kf.revenue_year})" if kf.revenue_year else kf.revenue),
        ("EBITDA:", f"{kf.ebitda} ({kf.ebitda_year})" if kf.ebitda_year else kf.ebitda),
        ("Management:", "\n".join(kf.management) if kf.management else ""),
        ("Employees:", kf.employees),
    ]

    first = True
    for label, value in fields:
        if not value:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(1)
        _add_run(p, label + " ", bold=True, color=DARK_NAVY, font_size=Pt(7))
        _add_run(p, value, font_size=Pt(7))


def _fill_bullets(shape_map: dict, shape_name: str, bullets: list[str]):
    """Fill a bullet-point text box."""
    shape = shape_map.get(shape_name)
    if not shape or not shape.has_text_frame:
        return

    tf = shape.text_frame
    _clear_text_frame(tf)

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        _add_run(p, f"• {bullet}", font_size=Pt(7))


def _fill_rationale(shape_map: dict, data: OnePagerData):
    """Fill the Investment Rationale (pros in green, cons in red) in column 4."""
    shape = shape_map.get("rationale_content")
    if not shape or not shape.has_text_frame:
        return

    tf = shape.text_frame
    _clear_text_frame(tf)

    first = True
    # Pros
    for pro in data.investment_rationale.pros:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(1)
        _add_run(p, f"+ {pro}", color=GREEN, font_size=Pt(7))

    # Cons
    for con in data.investment_rationale.cons:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(1)
        _add_run(p, f"– {con}", color=RED, font_size=Pt(7))


def _fill_criteria(shape_map: dict, data: OnePagerData):
    """Set the criteria indicator colors based on status."""
    criteria_map = {
        "criteria_ebitda_1m": data.investment_criteria.ebitda_1m,
        "criteria_dach": data.investment_criteria.dach,
        "criteria_ebitda_margin_10": data.investment_criteria.ebitda_margin_10,
        "criteria_majority_stake": data.investment_criteria.majority_stake,
        "criteria_digitization": data.investment_criteria.digitization,
        "criteria_asset_light": data.investment_criteria.asset_light,
        "criteria_buy_and_build": data.investment_criteria.buy_and_build,
        "criteria_esg": data.investment_criteria.esg,
        "criteria_market_fragmentation": data.investment_criteria.market_fragmentation,
        "criteria_acquisition_vertical": data.investment_criteria.acquisition_vertical,
        "criteria_acquisition_horizontal": data.investment_criteria.acquisition_horizontal,
        "criteria_acquisition_geographical": data.investment_criteria.acquisition_geographical,
    }

    status_colors = {
        CriterionStatus.FULFILLED: GREEN,
        CriterionStatus.QUESTIONS: GOLD,
        CriterionStatus.NOT_INTEREST: RED,
    }

    for prefix, status in criteria_map.items():
        indicator = shape_map.get(f"{prefix}_status")
        if indicator:
            color = status_colors.get(status, LIGHT_GRAY)
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = color


def _fill_status(shape_map: dict, data: OnePagerData):
    """Fill the Status section (column 4, bottom)."""
    shape = shape_map.get("status_content")
    if not shape or not shape.has_text_frame:
        return

    tf = shape.text_frame
    _clear_text_frame(tf)

    lines = [
        ("Status:", data.meta.status),
    ]

    # Add additional meta fields if present
    if data.meta.source:
        lines.insert(0, ("Source:", data.meta.source))
    if data.meta.im_received:
        lines.append(("IM received:", data.meta.im_received))
    if data.meta.loi_deadline:
        lines.append(("LOI Deadline:", data.meta.loi_deadline))

    for i, (label, value) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        _add_run(p, label + " ", bold=True, font_size=Pt(7))
        _add_run(p, value or "", font_size=Pt(7))


def _replace_shape_with_image(slide, shape_map: dict, shape_name: str,
                               img_bytes: bytes):
    """Replace a named shape with an image at the same position/size."""
    shape = shape_map.get(shape_name)
    if not shape:
        return

    left, top = shape.left, shape.top
    width, height = shape.width, shape.height

    # Remove the placeholder shape
    sp = shape._element
    sp.getparent().remove(sp)

    # Add image
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, left, top, width, height)


def generate_one_pager(data: OnePagerData, output_path: Optional[str] = None) -> bytes:
    """
    Generate a complete One-Pager PPTX from structured data.

    Uses the Constellation Capital AG template (derived from original
    hand-crafted templates) with the 4-column layout at 10" x 5.625".

    Args:
        data: Complete OnePagerData object
        output_path: Optional path to save the file

    Returns:
        PPTX file bytes
    """
    if not os.path.exists(TEMPLATE_PATH):
        raise FileNotFoundError(
            f"Template not found at {TEMPLATE_PATH}. "
            "Please ensure the template file exists."
        )

    # Copy template to temp file
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
    shutil.copy(TEMPLATE_PATH, tmp_path)

    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        shape_map = _build_shape_map(slide)

        # Fill text sections
        _fill_header(shape_map, data)
        _fill_key_facts(shape_map, data)
        _fill_bullets(shape_map, "description_content", data.description)
        _fill_bullets(shape_map, "portfolio_content", data.product_portfolio)
        _fill_rationale(shape_map, data)
        _fill_criteria(shape_map, data)
        _fill_status(shape_map, data)

        # Generate and embed charts
        if data.revenue_split.segments:
            donut_bytes = generate_revenue_donut(
                [s.model_dump() for s in data.revenue_split.segments],
                data.revenue_split.total,
            )
            _replace_shape_with_image(slide, shape_map, "revsplit_chart", donut_bytes)

        if data.financials.years:
            fin_bytes = generate_financials_chart(
                data.financials.years,
                data.financials.revenue,
                data.financials.ebitda,
                data.financials.ebitda_margin,
            )
            _replace_shape_with_image(slide, shape_map, "financials_chart", fin_bytes)

        # Save to bytes
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        result = buf.read()
    finally:
        os.unlink(tmp_path)

    # Optionally save to disk
    if output_path:
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        with open(output_path, "wb") as f:
            f.write(result)

    return result
