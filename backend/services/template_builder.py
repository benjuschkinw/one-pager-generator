"""
Template Builder: creates a clean base template from an original
Constellation Capital AG One-Pager PPTX.

The original templates (hand-crafted in PowerPoint with think-cell charts)
are stripped of company-specific content while preserving:
- Slide master, theme, colors, fonts
- 4-column layout structure (10" x 5.625")
- Section headers and separator lines
- Footer placeholders
- Investment criteria indicator shapes
- Color-coded legend

Named shapes are assigned for easy filling by the PPTX generator.
"""

import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _pos_key(shape):
    """Round shape position to 2 decimal places (inches)."""
    return (round(shape.left / 914400, 2), round(shape.top / 914400, 2))


def _size_key(shape):
    """Round shape size to 2 decimal places (inches)."""
    return (round(shape.width / 914400, 2), round(shape.height / 914400, 2))


def _clear_text_frame(tf):
    """Clear all text runs in a text frame."""
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""


# Shape renaming maps
RENAME_BY_POSITION = {
    (0.36, 0.24): "header_title",
    (0.36, 0.51): "header_tagline",
    (0.36, 0.92): "investment_thesis",
    (0.41, 1.50): "keyfacts_header",
    (2.45, 1.49): "description_header",
    (2.45, 3.39): "revsplit_header",
    (5.20, 1.49): "portfolio_header",
    (5.20, 3.39): "financials_header",
    (7.94, 1.49): "rationale_header",
    (7.94, 2.86): "criteria_header",
    (7.94, 3.91): "acquisition_header",
    (7.87, 4.65): "status_content",
}

CONTENT_RECTS = {
    ((0.35, 1.46), (1.87, 3.63)): "keyfacts_content",
    ((2.40, 1.46), (2.57, 1.76)): "description_content",
    ((5.14, 1.46), (2.57, 1.76)): "portfolio_content",
    ((7.86, 1.46), (1.78, 3.02)): "rationale_content",
    ((2.40, 3.34), (2.57, 1.75)): "revsplit_box",
    ((5.14, 3.34), (2.57, 1.75)): "financials_box",
}

CRITERIA_INDICATORS = {
    (7.95, 3.14): "criteria_ebitda_1m_status",
    (7.95, 3.30): "criteria_ebitda_margin_10_status",
    (7.95, 3.46): "criteria_digitization_status",
    (7.95, 3.62): "criteria_buy_and_build_status",
    (7.95, 3.79): "criteria_market_fragmentation_status",
    (8.94, 3.14): "criteria_dach_status",
    (8.94, 3.30): "criteria_majority_stake_status",
    (8.94, 3.46): "criteria_asset_light_status",
    (8.94, 3.62): "criteria_esg_status",
    (7.95, 4.19): "criteria_acquisition_vertical_status",
    (8.94, 4.19): "criteria_acquisition_horizontal_status",
    (7.95, 4.33): "criteria_acquisition_geographical_status",
}

CONTENT_SHAPES_TO_CLEAR = {
    "keyfacts_content", "description_content", "portfolio_content",
    "rationale_content", "status_content", "investment_thesis",
    "header_tagline",
}


def build_template(source_pptx: str, output_path: str) -> str:
    """
    Build a clean base template from an original Constellation Capital AG
    One-Pager PPTX file.

    Args:
        source_pptx: Path to the original hand-crafted PPTX
        output_path: Path to save the clean template

    Returns:
        The output path
    """
    prs = Presentation(source_pptx)
    slide = prs.slides[0]

    # Remove extra slides (keep only slide 1)
    from pptx.oxml.ns import qn
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[-1].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

    # Phase 1: Rename content shapes
    for shape in slide.shapes:
        pk = _pos_key(shape)
        sk = _size_key(shape)

        if pk in RENAME_BY_POSITION:
            shape.name = RENAME_BY_POSITION[pk]

        for (pos, sz), new_name in CONTENT_RECTS.items():
            if pk == pos and sk == sz:
                shape.name = new_name

    # Phase 2: Rename criteria indicators (0.08x0.08 rectangles)
    for shape in slide.shapes:
        pk = _pos_key(shape)
        sk = _size_key(shape)
        if pk in CRITERIA_INDICATORS and sk == (0.08, 0.08):
            shape.name = CRITERIA_INDICATORS[pk]

    # Phase 3: Clear text content from data shapes
    for shape in slide.shapes:
        if shape.name in CONTENT_SHAPES_TO_CLEAR and shape.has_text_frame:
            _clear_text_frame(shape.text_frame)

        if shape.name == "header_title" and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "One Pager" not in run.text:
                        run.text = ""

    # Phase 4: Remove charts, images, OLE objects, and chart annotations
    shapes_to_remove = []
    chart_replacements = []

    for shape in slide.shapes:
        pk = _pos_key(shape)

        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_replacements.append({
                'left': shape.left, 'top': shape.top,
                'width': shape.width, 'height': shape.height,
                'name': 'revsplit_chart' if shape.left < 4 * 914400 else 'financials_chart'
            })
            shapes_to_remove.append(shape)

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_to_remove.append(shape)

        elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            shapes_to_remove.append(shape)

        elif shape.name.startswith('Oval'):
            shapes_to_remove.append(shape)

        elif shape.name.startswith('Straight Connector'):
            shapes_to_remove.append(shape)

        # Financial chart annotations (small labeled rectangles in chart areas)
        elif shape.name.startswith('Rectangle') and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if pk[0] >= 5.0 and pk[1] >= 3.7:
                if text and (text.endswith('A') or text.endswith('P') or text.endswith('E') or
                            '%' in text or text.replace('.', '').replace('-', '').replace('(', '').replace(')', '').isdigit()):
                    shapes_to_remove.append(shape)

        # Revenue split legend items
        elif shape.name.startswith('Rectangle') and pk[0] >= 3.8 and pk[0] < 5.0 and pk[1] >= 3.7:
            shapes_to_remove.append(shape)

    # EUR label and total label
    for shape in slide.shapes:
        pk = _pos_key(shape)
        if pk == (5.17, 3.61) and shape.has_text_frame and 'EUR' in shape.text_frame.text:
            shapes_to_remove.append(shape)
        if shape.name == 'Rectangle 1093':
            shapes_to_remove.append(shape)

    # Remove table and D&A text box
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            shapes_to_remove.append(shape)
        if shape.name == 'TextBox 1397':
            shapes_to_remove.append(shape)

    # Execute removals
    removed = set()
    for shape in shapes_to_remove:
        sp = shape._element
        if sp not in removed and sp.getparent() is not None:
            sp.getparent().remove(sp)
            removed.add(sp)

    # Add chart placeholder shapes
    for repl in chart_replacements:
        txBox = slide.shapes.add_textbox(repl['left'], repl['top'], repl['width'], repl['height'])
        txBox.name = repl['name']

    # Save
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python template_builder.py <source_pptx> [output_path]")
        print("  source_pptx: Path to an original Constellation Capital PPTX")
        print("  output_path: Where to save the clean template (default: ../template/one_pager_template.pptx)")
        sys.exit(1)

    source = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else os.path.join(
        os.path.dirname(__file__), "..", "template", "one_pager_template.pptx"
    )
    result = build_template(source, output)
    print(f"Template saved to: {result}")
